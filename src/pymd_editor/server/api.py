"""
PyMD Editor — FastAPI REST API
Provides all backend logic for the web frontend.
Security: server binds to 127.0.0.1 only; all file paths are validated
against BASE_DIR to prevent path traversal.
"""
from __future__ import annotations

import io
import os
import tempfile
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, HTTPException, Query, BackgroundTasks, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel

# ---------------------------------------------------------------------------
# Bootstrap import path so this module can be run standalone or imported
# ---------------------------------------------------------------------------
import sys
_src = Path(__file__).parent.parent.parent
if str(_src) not in sys.path:
    sys.path.insert(0, str(_src))

from pymd_editor.renderer import MarkdownRenderer
from pymd_editor.exporter import WordExporter
from pymd_editor.config import APP_VERSION

# ---------------------------------------------------------------------------
# App
# ---------------------------------------------------------------------------

app = FastAPI(title="PyMD Editor", version=APP_VERSION, docs_url="/api/docs")

def _cors_origins() -> list[str]:
    defaults = [
        "http://localhost:8765",
        "http://127.0.0.1:8765",
    ]
    extra = [
        origin.strip()
        for origin in os.getenv("PYMD_CORS_ALLOW_ORIGINS", "").split(",")
        if origin.strip()
    ]
    return defaults + extra


app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_origins(),
    allow_methods=["GET", "POST", "DELETE", "OPTIONS"],
    allow_headers=["*"],
)

_renderer = MarkdownRenderer()
_word_exporter = WordExporter()

# Base directory for ALL file operations — overridden by serve.py at startup
BASE_DIR: Path = Path.home()
DEPLOYMENT_MODE = os.getenv("PYMD_DEPLOYMENT_MODE", "local")
PUBLIC_BASE_URL = os.getenv("PYMD_PUBLIC_BASE_URL", "").strip()


def _feature_flags() -> dict[str, bool]:
    return {
        "markdown_render": True,
        "markdown_export_word": True,
        "folder_browse": True,
        "pdf_info": True,
        "pdf_extract": True,
        "pdf_merge": True,
        "pdf_insert": True,
        "pdf_to_word": True,
        "pdf_to_excel": True,
        "pdf_to_ppt": True,
    }


def _convert_pdf_path_to_word(pdf_path: Path, out_path: Path) -> None:
    from pdf2docx import Converter

    cv = Converter(str(pdf_path))
    try:
        cv.convert(str(out_path), start=0, end=None)
    finally:
        cv.close()


def _convert_pdf_path_to_excel(pdf_path: Path, out_path: Path) -> None:
    import fitz, openpyxl

    doc = fitz.open(str(pdf_path))
    try:
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        for pg_num, page in enumerate(doc, 1):
            ws = wb.create_sheet(title=f"Page {pg_num}")
            blocks = sorted(page.get_text("blocks"), key=lambda b: (round(b[1] / 10) * 10, b[0]))
            row = 1
            for blk in blocks:
                if blk[6] != 0:
                    continue
                for line in blk[4].strip().splitlines():
                    cells = [c.strip() for c in line.split("\t")] if "\t" in line else [line]
                    for col, val in enumerate(cells, 1):
                        ws.cell(row=row, column=col, value=val)
                    row += 1
        wb.save(str(out_path))
    finally:
        doc.close()


def _convert_pdf_path_to_ppt(pdf_path: Path, out_path: Path) -> None:
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(str(pdf_path))
    try:
        prs = Presentation()
        prs.slide_width = Emu(13_333_333)
        prs.slide_height = Emu(7_500_000)
        blank = prs.slide_layouts[6]

        for page in doc:
            slide = prs.slides.add_slide(blank)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_bytes = pix.tobytes("png")

            page_rect = page.rect
            page_ratio = page_rect.width / page_rect.height if page_rect.height else 1.0
            slide_ratio = prs.slide_width / prs.slide_height

            if page_ratio >= slide_ratio:
                width = prs.slide_width
                height = int(width / page_ratio)
                left = 0
                top = int((prs.slide_height - height) / 2)
            else:
                height = prs.slide_height
                width = int(height * page_ratio)
                top = 0
                left = int((prs.slide_width - width) / 2)

            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)

        if len(prs.slides) > 1:
            first_slide = prs.slides._sldIdLst[0]
            prs.slides._sldIdLst.remove(first_slide)

        prs.save(str(out_path))
    finally:
        doc.close()


def _tmp_named(suffix: str) -> Path:
    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    path = Path(tmp.name)
    tmp.close()
    return path


# ---------------------------------------------------------------------------
# Security helper
# ---------------------------------------------------------------------------

def _safe_path(rel_or_abs: str) -> Path:
    """Resolve *rel_or_abs* and verify it stays within BASE_DIR.

    Raises HTTPException 400 for invalid paths and 403 for path-traversal
    attempts.
    """
    try:
        p = Path(rel_or_abs)
        if not p.is_absolute():
            p = BASE_DIR / p
        p = p.resolve()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")

    base = BASE_DIR.resolve()
    try:
        p.relative_to(base)
    except ValueError:
        raise HTTPException(status_code=403, detail="Path outside allowed directory")

    return p


# ---------------------------------------------------------------------------
# Request / response models
# ---------------------------------------------------------------------------

class RenderRequest(BaseModel):
    markdown: str
    dark: bool = False
    base_path: Optional[str] = None


class FileWriteRequest(BaseModel):
    path: str
    content: str


class ExportWordRequest(BaseModel):
    markdown: str
    filename: str = "document.docx"


# ---------------------------------------------------------------------------
# Routes — Health
# ---------------------------------------------------------------------------

@app.get("/api/health")
def health():
    return {
        "status": "ok",
        "version": APP_VERSION,
        "base_dir": str(BASE_DIR),
        "deployment_mode": DEPLOYMENT_MODE,
        "public_base_url": PUBLIC_BASE_URL,
        "features": _feature_flags(),
    }


# ---------------------------------------------------------------------------
# Routes — Render
# ---------------------------------------------------------------------------

@app.post("/api/render")
def render(req: RenderRequest):
    """Convert markdown text to full HTML (with CSS + MathJax)."""
    base_path: Optional[str] = None
    if req.base_path:
        # Validate; silently ignore if outside BASE_DIR
        try:
            bp = _safe_path(req.base_path)
            base_path = str(bp)
        except HTTPException:
            pass

    html = _renderer.to_html(req.markdown, dark=req.dark, base_path=base_path)
    return {"html": html}


# ---------------------------------------------------------------------------
# Routes — File operations
# ---------------------------------------------------------------------------

@app.get("/api/files")
def list_files(
    dir: str = Query(default=""),
    types: str = Query(default="md"),          # comma-separated: md,pdf
):
    """List files in a directory. types=md,pdf to include PDFs."""
    target = _safe_path(dir) if dir else BASE_DIR
    if not target.is_dir():
        raise HTTPException(status_code=400, detail="Not a directory")

    allowed_exts = set()
    for t in types.split(","):
        t = t.strip().lower().lstrip(".")
        if t in ("md", "pdf"):
            allowed_exts.add("." + t)

    items = []
    try:
        for item in sorted(target.iterdir()):
            rel = str(item.relative_to(BASE_DIR))
            if item.is_dir() and not item.name.startswith("."):
                items.append({"name": item.name, "path": rel, "type": "dir"})
            elif item.is_file() and item.suffix.lower() in allowed_exts:
                stat = item.stat()
                items.append({
                    "name": item.name,
                    "path": rel,
                    "type": "file",
                    "ext":  item.suffix.lower().lstrip("."),
                    "size": stat.st_size,
                    "modified": stat.st_mtime,
                })
    except PermissionError:
        raise HTTPException(status_code=403, detail="Permission denied")

    current_rel = str(target.relative_to(BASE_DIR)) if target != BASE_DIR else ""
    return {"files": items, "dir": current_rel}


@app.get("/api/file")
def read_file(path: str = Query(...)):
    p = _safe_path(path)
    if not p.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    if p.suffix.lower() != ".md":
        raise HTTPException(status_code=400, detail="Only .md files are supported")

    content = p.read_text(encoding="utf-8")
    rel = str(p.relative_to(BASE_DIR))
    return {"content": content, "path": rel, "name": p.name}


@app.post("/api/file")
def write_file(req: FileWriteRequest):
    p = _safe_path(req.path)
    if p.suffix.lower() != ".md":
        raise HTTPException(status_code=400, detail="Only .md files are supported")

    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(req.content, encoding="utf-8")
    return {"ok": True, "path": str(p.relative_to(BASE_DIR))}


@app.delete("/api/file")
def delete_file(path: str = Query(...)):
    p = _safe_path(path)
    if not p.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    if p.suffix.lower() != ".md":
        raise HTTPException(status_code=400, detail="Only .md files are supported")

    p.unlink()
    return {"ok": True}


# ---------------------------------------------------------------------------
# Routes — Export
# ---------------------------------------------------------------------------

@app.post("/api/export/word")
def export_word(req: ExportWordRequest, background_tasks: BackgroundTasks):
    """Export markdown as a .docx file (streamed download)."""
    # Strip any path components from the requested filename
    fname = Path(req.filename).name
    if not fname.lower().endswith(".docx"):
        fname += ".docx"

    tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    tmp_path = Path(tmp.name)
    tmp.close()

    try:
        _word_exporter.export(req.markdown, tmp_path)
    except Exception as exc:
        tmp_path.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=str(exc))

    # Clean up temp file after FastAPI finishes streaming the response
    background_tasks.add_task(tmp_path.unlink, True)

    return FileResponse(
        path=str(tmp_path),
        filename=fname,
        media_type=(
            "application/vnd.openxmlformats-officedocument"
            ".wordprocessingml.document"
        ),
    )


# ===========================================================================
# PDF helpers & endpoints
# ===========================================================================

def _require_pdf(p: Path) -> None:
    if p.suffix.lower() != ".pdf":
        raise HTTPException(400, "Only .pdf files are supported")
    if not p.is_file():
        raise HTTPException(404, "File not found")


def _parse_page_ranges(spec: str, total: int) -> list[int]:
    """'1-3, 5, 7' → [0,1,2,4,6] (0-indexed, clamped to valid range)."""
    result: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            lo_s, hi_s = part.split("-", 1)
            for i in range(int(lo_s.strip()) - 1, int(hi_s.strip())):
                if 0 <= i < total:
                    result.add(i)
        elif part.isdigit():
            i = int(part) - 1
            if 0 <= i < total:
                result.add(i)
    return sorted(result)


def _pdf_tmp() -> Path:
    t = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    p = Path(t.name); t.close(); return p


# ── PDF request models ────────────────────────────────────────────────────────

class PDFExtractRequest(BaseModel):
    path: str
    pages: str           # "1-3, 5, 7"
    filename: str = ""

class PDFMergeRequest(BaseModel):
    paths: list[str]
    filename: str = "merged.pdf"

class PDFInsertRequest(BaseModel):
    base: str
    insert: str
    after_page: int      # 1-indexed; 0 = before page 1
    filename: str = ""

class PDFToWordRequest(BaseModel):
    path: str
    filename: str = ""

class PDFToExcelRequest(BaseModel):
    path: str
    filename: str = ""


# ── PDF stream (in-browser viewer) ──────────────────────────────────────────

@app.get("/api/pdf/stream")
def pdf_stream(path: str = Query(...)):
    """Serve a PDF inline so the browser can render it in an <iframe>."""
    p = _safe_path(path)
    _require_pdf(p)
    return FileResponse(
        str(p),
        media_type="application/pdf",
        headers={"Content-Disposition": f'inline; filename="{p.name}"'},
    )


# ── PDF info & thumbnail ──────────────────────────────────────────────────────

@app.get("/api/pdf/info")
def pdf_info(path: str = Query(...)):
    p = _safe_path(path)
    _require_pdf(p)
    try:
        import fitz
        doc  = fitz.open(str(p))
        meta = doc.metadata
        result = {
            "pages":       len(doc),
            "title":       meta.get("title", ""),
            "author":      meta.get("author", ""),
            "size_bytes":  p.stat().st_size,
            "name":        p.name,
            "path":        str(p.relative_to(BASE_DIR)),
        }
        doc.close()
        return result
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.get("/api/pdf/thumb")
def pdf_thumb(path: str = Query(...), page: int = Query(default=1)):
    """Return a PNG thumbnail of a PDF page."""
    from fastapi.responses import Response as _Resp
    p = _safe_path(path)
    _require_pdf(p)
    try:
        import fitz
        doc   = fitz.open(str(p))
        total = len(doc)
        if page < 1 or page > total:
            raise HTTPException(400, f"Page must be 1–{total}")
        pix = doc.load_page(page - 1).get_pixmap(matrix=fitz.Matrix(0.8, 0.8))
        data = pix.tobytes("png")
        doc.close()
        return _Resp(content=data, media_type="image/png")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


# ── PDF operations ────────────────────────────────────────────────────────────

@app.post("/api/pdf/extract")
def pdf_extract(req: PDFExtractRequest, bg: BackgroundTasks):
    p = _safe_path(req.path);  _require_pdf(p)
    try:
        import fitz
        doc   = fitz.open(str(p))
        pages = _parse_page_ranges(req.pages, len(doc))
        if not pages:
            raise HTTPException(400, "No valid pages in that range")
        out = fitz.open()
        for i in pages:
            out.insert_pdf(doc, from_page=i, to_page=i)
        tmp = _pdf_tmp()
        out.save(str(tmp));  out.close();  doc.close()
        fname = Path(req.filename or (p.stem + "_extract.pdf")).name
        bg.add_task(tmp.unlink, True)
        return FileResponse(str(tmp), filename=fname, media_type="application/pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/pdf/merge")
def pdf_merge(req: PDFMergeRequest, bg: BackgroundTasks):
    if len(req.paths) < 2:
        raise HTTPException(400, "Need at least 2 PDFs to merge")
    paths = [_safe_path(p) for p in req.paths]
    for p in paths:
        _require_pdf(p)
    try:
        import fitz
        merged = fitz.open()
        for p in paths:
            d = fitz.open(str(p));  merged.insert_pdf(d);  d.close()
        tmp = _pdf_tmp()
        merged.save(str(tmp));  merged.close()
        fname = Path(req.filename or "merged.pdf").name
        bg.add_task(tmp.unlink, True)
        return FileResponse(str(tmp), filename=fname, media_type="application/pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/pdf/insert")
def pdf_insert(req: PDFInsertRequest, bg: BackgroundTasks):
    base_p = _safe_path(req.base);    _require_pdf(base_p)
    ins_p  = _safe_path(req.insert);  _require_pdf(ins_p)
    try:
        import fitz
        base = fitz.open(str(base_p))
        ins  = fitz.open(str(ins_p))
        at   = max(0, min(req.after_page, len(base)))
        base.insert_pdf(ins, start_at=at);  ins.close()
        tmp  = _pdf_tmp()
        base.save(str(tmp));  base.close()
        fname = Path(req.filename or (base_p.stem + "_inserted.pdf")).name
        bg.add_task(tmp.unlink, True)
        return FileResponse(str(tmp), filename=fname, media_type="application/pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/pdf/to-word")
def pdf_to_word(req: PDFToWordRequest, bg: BackgroundTasks):
    p = _safe_path(req.path);  _require_pdf(p)
    try:
        tmp_path = _tmp_named(".docx")
        _convert_pdf_path_to_word(p, tmp_path)
        fname = Path(req.filename or (p.stem + ".docx")).name
        bg.add_task(tmp_path.unlink, True)
        return FileResponse(
            str(tmp_path), filename=fname,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/pdf/to-excel")
def pdf_to_excel(req: PDFToExcelRequest, bg: BackgroundTasks):
    p = _safe_path(req.path);  _require_pdf(p)
    try:
        tmp_path = _tmp_named(".xlsx")
        _convert_pdf_path_to_excel(p, tmp_path)
        fname = Path(req.filename or (p.stem + ".xlsx")).name
        bg.add_task(tmp_path.unlink, True)
        return FileResponse(
            str(tmp_path), filename=fname,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


class PDFToPptRequest(BaseModel):
    path: str
    filename: Optional[str] = None


@app.post("/api/pdf/to-ppt")
def pdf_to_ppt(req: PDFToPptRequest, bg: BackgroundTasks):
    p = _safe_path(req.path);  _require_pdf(p)
    try:
        tmp_path = _tmp_named(".pptx")
        _convert_pdf_path_to_ppt(p, tmp_path)
        fname = Path(req.filename or (p.stem + ".pptx")).name
        bg.add_task(tmp_path.unlink, True)
        return FileResponse(
            str(tmp_path), filename=fname,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


def _uploaded_pdf_to_temp(upload: UploadFile) -> tuple[Path, str]:
    suffix = Path(upload.filename or "upload.pdf").suffix or ".pdf"
    tmp_path = _tmp_named(suffix)
    with tmp_path.open("wb") as f:
        while True:
            chunk = upload.file.read(1024 * 1024)
            if not chunk:
                break
            f.write(chunk)
    return tmp_path, Path(upload.filename or "document.pdf").stem


@app.post("/api/pdf/upload/to-word")
def pdf_upload_to_word(
    bg: BackgroundTasks,
    file: UploadFile = File(...),
    filename: Optional[str] = Form(default=None),
):
    src_path, stem = _uploaded_pdf_to_temp(file)
    try:
        _require_pdf(src_path)
        out_path = _tmp_named(".docx")
        _convert_pdf_path_to_word(src_path, out_path)
        out_name = Path(filename or (stem + ".docx")).name
        bg.add_task(src_path.unlink, True)
        bg.add_task(out_path.unlink, True)
        return FileResponse(
            str(out_path), filename=out_name,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    except Exception:
        src_path.unlink(missing_ok=True)
        raise


@app.post("/api/pdf/upload/to-excel")
def pdf_upload_to_excel(
    bg: BackgroundTasks,
    file: UploadFile = File(...),
    filename: Optional[str] = Form(default=None),
):
    src_path, stem = _uploaded_pdf_to_temp(file)
    try:
        _require_pdf(src_path)
        out_path = _tmp_named(".xlsx")
        _convert_pdf_path_to_excel(src_path, out_path)
        out_name = Path(filename or (stem + ".xlsx")).name
        bg.add_task(src_path.unlink, True)
        bg.add_task(out_path.unlink, True)
        return FileResponse(
            str(out_path), filename=out_name,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
    except Exception:
        src_path.unlink(missing_ok=True)
        raise


@app.post("/api/pdf/upload/to-ppt")
def pdf_upload_to_ppt(
    bg: BackgroundTasks,
    file: UploadFile = File(...),
    filename: Optional[str] = Form(default=None),
):
    src_path, stem = _uploaded_pdf_to_temp(file)
    try:
        _require_pdf(src_path)
        out_path = _tmp_named(".pptx")
        _convert_pdf_path_to_ppt(src_path, out_path)
        out_name = Path(filename or (stem + ".pptx")).name
        bg.add_task(src_path.unlink, True)
        bg.add_task(out_path.unlink, True)
        return FileResponse(
            str(out_path), filename=out_name,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception:
        src_path.unlink(missing_ok=True)
        raise


# ===========================================================================
# File management endpoints
# ===========================================================================

class FileRenameRequest(BaseModel):
    path: str
    new_name: str    # bare filename only

class FolderCreateRequest(BaseModel):
    path: str        # relative path for new folder

class FileMoveRequest(BaseModel):
    path: str
    dest: str
    do_copy: bool = False   # renamed from 'copy' to avoid Pydantic BaseModel clash


@app.post("/api/file/rename")
def rename_file(req: FileRenameRequest):
    src = _safe_path(req.path)
    if not src.is_file():
        raise HTTPException(404, "File not found")
    new_name = Path(req.new_name).name
    if not new_name or new_name in (".", ".."):
        raise HTTPException(400, "Invalid filename")
    dst = (src.parent / new_name).resolve()
    try:
        dst.relative_to(BASE_DIR.resolve())
    except ValueError:
        raise HTTPException(403, "Destination outside allowed directory")
    if dst.exists():
        raise HTTPException(409, "A file with that name already exists")
    src.rename(dst)
    return {"ok": True, "path": str(dst.relative_to(BASE_DIR))}


@app.post("/api/folder")
def create_folder(req: FolderCreateRequest):
    p = _safe_path(req.path)
    if p.exists():
        raise HTTPException(409, "Already exists")
    p.mkdir(parents=True)
    return {"ok": True, "path": str(p.relative_to(BASE_DIR))}


@app.post("/api/file/move")
def move_file(req: FileMoveRequest):
    import shutil
    src  = _safe_path(req.path)
    dest = _safe_path(req.dest)
    if not src.is_file():
        raise HTTPException(404, "Source not found")
    if dest.exists() and not req.do_copy:
        raise HTTPException(409, "Destination already exists")
    dest.parent.mkdir(parents=True, exist_ok=True)
    if req.do_copy:
        shutil.copy2(str(src), str(dest))
    else:
        src.rename(dest)
    return {"ok": True, "path": str(dest.relative_to(BASE_DIR))}
